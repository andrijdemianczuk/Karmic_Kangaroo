
#Import types
from typing import Annotated, Any, Generator, Optional, Sequence, TypedDict, Union, TypedDict, List, Dict

import os

#MLFlow libs
import mlflow
from mlflow.pyfunc import ResponsesAgent
from mlflow.types.responses import (
    ResponsesAgentRequest,
    ResponsesAgentResponse,
    ResponsesAgentStreamEvent,
    output_to_responses_items_stream,
    to_chat_completions_input,
)

#Databricks MLFlow flavour of LangChain for Databricks interop
from databricks_langchain import (
    ChatDatabricks,
    UCFunctionToolkit,
    VectorSearchRetrieverTool,
)

#Standard Langchain libs
from langchain_core.messages import AIMessage, AIMessageChunk, AnyMessage
from langchain_core.runnables import RunnableConfig, RunnableLambda
from langchain_core.tools import BaseTool
from langchain.tools import tool
from langchain_openai import ChatOpenAI

#Langgraph for orchestration
from langgraph.graph import END, StateGraph
from langgraph.graph.message import add_messages
from langgraph.prebuilt.tool_node import ToolNode

#Extended Unity Catalog AI capabilities (don't use the DatabricksFunctionClient() from databricks_langchain)
# from unitycatalog.ai.core.databricks import DatabricksFunctionClient

#Standard protocols
import json

#Microsoft Powerpoint libraries
from pptx import Presentation

#Microsoft Word libraries
from docx import Document
from docx import Document as load_docx  # constructor function
from docx.document import Document as DocxDocument  # actual document type
from docx.table import Table as DocxTable, _Cell
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl


#Create a struct to keep track of the evaluation agent
class EvalState(TypedDict, total=False):
    file_name: str                      # e.g. "incident_123.pptx"
    user_prompt: str                    # any prompt / focus area text
    pptx_chunks: List[str]              # parsed slide texts
    docx_chunks: List[str]              # parsed paragraph texts
    corrective_actions_raw: str         # raw LLM output from gen agent
    corrective_actions: List[Dict]      # parsed actions from gen agent
    evaluated_actions_raw: str          # raw LLM output from eval agent
    evaluated_actions: List[Dict]       # top-k ranked actions


#Define the endpoint to use for the agent foundation and system prompt
LLM_ENDPOINT_NAME = "databricks-claude-3-7-sonnet"
llm = ChatDatabricks(endpoint=LLM_ENDPOINT_NAME)


#Change this to wherever our user app uploads documents
UPLOAD_VOLUME = "/Volumes/ademianczuk/suncor_ehs/data/uploads"

#Break up the document into manageable blocks
def _iter_block_items(parent):
    """
    Yield paragraphs and tables in document order.

    This uses the underlying oxml tree to preserve the order of paragraphs and tables.
    Works for a docx.document.Document and table cells (_Cell).
    """
    if isinstance(parent, DocxDocument):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError(f"Unsupported parent type for iter_block_items: {type(parent)}")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield DocxParagraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)


#Break up and parse the tables. This is what was missing in the last version of this tool
def _table_to_text(table: DocxTable, table_index: int) -> str:
    """
    Convert a python-docx Table into a readable text block.

    We treat the first row as headers (if present) and then produce lines like:
      - Header1: val1; Header2: val2; ...
    """
    rows = list(table.rows)
    if not rows:
        return f"TABLE {table_index}: (empty)"

    #First row as headers
    header_cells = [cell.text.strip() for cell in rows[0].cells]
    headers = [h if h else f"Column {i+1}" for i, h in enumerate(header_cells)]

    lines: List[str] = []

    #Data rows
    for row in rows[1:]:
        cells = [cell.text.strip() for cell in row.cells]
        pairs = []
        for i, value in enumerate(cells):
            if not value:
                continue
            header = headers[i] if i < len(headers) else f"Column {i+1}"
            pairs.append(f"{header}: {value}")
        if pairs:
            lines.append("- " + "; ".join(pairs))

    if not lines:
        return f"TABLE {table_index}: (no data rows)"

    return f"TABLE {table_index}:\n" + "\n".join(lines)


#The actual tool definition. The decorator is imported from LangChain tools and provides the wrapper context necessary to declare code-as-tools.
@tool("parse_docx")
def parse_docx_tool(file_name: str) -> List[str]:
    """
    Parse a Word (.docx) file stored in the UC volume into text chunks.

    - Includes BOTH paragraph text and table contents.
    - Tables are converted into readable text blocks (TABLE n: ...).
    - Use this tool when the file has a .docx extension.
    - Files are expected under /Volumes/ademianczuk/suncor_ehs/data/uploads.
    """
    #Copy from UC volume (dbfs) to local FS so python-docx can open it

    ################
    #Needs to get refactored not to use dbutils
    ################

    # dbfs_path = f"dbfs:{UPLOAD_VOLUME}/{file_name}"
    # local_path = "/tmp/input.docx"
    upload_path = f"{UPLOAD_VOLUME}/{file_name}"

    # dbutils.fs.cp(dbfs_path, f"file:{local_path}", True)

    #Load DOCX (this returns an instance of DocxDocument - This is what was broken before :p)
    # doc = load_docx(local_path)
    doc = load_docx(upload_path)

    #Collect text blocks in document order (paragraphs + tables)
    blocks: List[str] = []
    table_counter = 0

    #Iterate through each block - the LM handles positional undersanding of table chunks (vertical & horizontal)
    for block in _iter_block_items(doc):
        if isinstance(block, DocxParagraph):
            text = block.text.strip()
            if text:
                blocks.append(text)
        elif isinstance(block, DocxTable):
            table_counter += 1
            table_text = _table_to_text(block, table_counter)
            blocks.append(table_text)

    #Chunk blocks into ~1200-character chunks
    chunks: List[str] = []
    current = ""
    max_chars = 1200

    for b in blocks:
        #+2 for newline spacing
        if len(current) + len(b) + 2 > max_chars:
            if current:
                chunks.append(current.strip())
            current = b
        else:
            current = f"{current}\n\n{b}" if current else b

    if current:
        chunks.append(current.strip())

    return chunks


@tool("parse_pptx")
def parse_pptx_tool(file_name: str) -> list[str]:
    """
    Parse a PPTX file stored in the UC Volume into slide-level text chunks. Use this to parse a Microsoft Powerpoint (.pptx) file.
    
    Args:
        file_name: The PPTX filename (e.g. 'incident_123.pptx'), assumed to live under
                   /Volumes/ademianczuk/suncor_ehs/data/uploads.

    Returns:
        A list of text chunks (strings). Each chunk is roughly slide-level; 
        small slides may be merged into bigger chunks to keep context.
    """
    # dbfs_path = f"dbfs:{UPLOAD_VOLUME}/{file_name}"
    # local_path = "/tmp/input.pptx"
    upload_path = f"{UPLOAD_VOLUME}/{file_name}"
    # dbutils.fs.cp(dbfs_path, f"file:{local_path}", True)

    prs = Presentation(upload_path)
    chunks: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)
        if texts:
            #each slide = one chunk for now
            chunks.append(f"Slide {slide_idx}:\n" + "\n".join(texts))

    return chunks

class AgentState(TypedDict):
    messages: Annotated[Sequence[AnyMessage], add_messages]
    custom_inputs: Optional[dict[str, Any]]
    custom_outputs: Optional[dict[str, Any]]


def create_tool_calling_agent(
    model: ChatDatabricks,
    tools: Union[ToolNode, Sequence[BaseTool]],
    system_prompt: Optional[str] = None,
):
    model = model.bind_tools(tools)

    # Define the function that determines which node to go to
    def should_continue(state: AgentState):
        messages = state["messages"]
        last_message = messages[-1]
        # If there are function calls, continue. else, end
        if isinstance(last_message, AIMessage) and last_message.tool_calls:
            return "continue"
        else:
            return "end"

    if system_prompt:
        preprocessor = RunnableLambda(
            lambda state: [{"role": "system", "content": system_prompt}] + state["messages"]
        )
    else:
        preprocessor = RunnableLambda(lambda state: state["messages"])
    model_runnable = preprocessor | model

    def call_model(
        state: AgentState,
        config: RunnableConfig,
    ):
        response = model_runnable.invoke(state, config)

        return {"messages": [response]}

    workflow = StateGraph(AgentState)

    workflow.add_node("agent", RunnableLambda(call_model))
    workflow.add_node("tools", ToolNode(tools))

    workflow.set_entry_point("agent")
    workflow.add_conditional_edges(
        "agent",
        should_continue,
        {
            "continue": "tools",
            "end": END,
        },
    )
    workflow.add_edge("tools", "agent")

    return workflow.compile()
    
    

CORRECTIVE_ACTIONS_INSTRUCTIONS = """
You are a corporate Enterprise Health and Safety (EHS) expert specializing in corrective actions.

With the following guidance, define and create corrective actions:

1. Before reading the remaining instructions below, you must observe all uploaded documents and ensure that you have a full and comprehensive understanding of them. You must ensure that you are able to extrapolate from the provided resources, directly following the framework, definitions, and standards provided in doing so. Make sure that you are looking for pattern recognition and defining areas of interest that are related to the goal trying to be achieved.

2. You will be provided with an incident investigation report. This may be presented in a PDF, Word, or PowerPoint format and you must be ready to read, interpret, and understand the content provided by the user. Your goal is to flag the report for instances of 'negative' reasoning. The definitions for these terms are provided in the uploaded document labelled 'Consolidated Context Format.docx'. Additionally, you must also identify the presence of counterfactuals in the report provided. In carrying out these defined tasks, you must also determine the difference between negative and causal reasoning as you will need this for analysis.

3. When finding instances of negative reasoning, showcase how you would define and extract instances of negative reasoning for future inquiries. You must determine what factors constitute the usage of negative reasoning in any given incident investigations report.

4. You must structure the output as follows. I want you to have three headers: 'Relationship', 'Flagging', and 'Corrective Actions'.

5. Under 'Relationship', reassess the prompt provided by the user, as well as the uploaded document labelled 'Consolidated Context Format.docx'. Review the incident investigation report and specifically identify any terms and definitions from the uploaded documents that are present. Identify specific excerpts of the investigation report that contain any terms and definitions. Explain how they are present. Do not be vague and try to find the presence of as many terms and definitions as you can with sufficient detail. Do not combine any definitions.

6. Under 'Flagging', specify all instances of negative reasoning in the report, as well as the presence of counterfactuals and any logical errors. Put these into a table with four columns, 'Identification of Negative Reasoning / Counterfactual' (Identify whether you are assessing an instance of only 'Negative Reasoning' or a 'Counterfactual', or combine both together if a certain excerpt has both an instance of negative reasoning and a counterfactual), 'Identification Reasoning' (How did you identify a given instance of negative reasoning or the usage of a counterfactual? What factors revealed their presence?), 'Definitions Present' (Which of the definitions above are present in the excerpt extracted?), and 'Original Statement' (The original statement being analyzed). Ensure that you are specific and detailed in every column.

7. Under 'Corrective Actions',

I want you to list the causes and actions from the investigation report in a sentence format above the table. I then want you to create a table that is populated with every single cause and action in the report from the initial list you made, in the order specified under the forthcoming column called 'Quality of Action'. Before creating the table with the following columns, ensure the initial lists' actions are ranked according to the 'Quality of Action' column below before the table is created from the list. The table should be structured as follows with the following logic:
    - 'Tally' (Just increments and counts every row)
    - 'Action' (the corresponding corrective action)
    - 'Cause' (The initial cause that the action is a result of)
    - 'OEMS Process' (the associated and most applicable OEMS Process defined in the document labelled 'OEMS Process Descriptions.docx' to the cause)
    -  'Related OEMS Process' (to showcase all other most applicable OEMS Processes (more than one) from the document labelled 'OEMS Process Descriptions.docx')
    - 'Hierarchy of Controls' (which determines the control hierarchy level the corrective action best embodies, defined in the document labelled 'Consolidated Context Format.docx' under the section titled 'Hierarchy of Controls - Corrective Actions')
    - 'Quality of Action (which ranks the action based on what control hierarchy level it is at (create a scale where Elimination would be a 5 (Most Effective), Substitution would be a 4 (Highly Effective), Engineering Controls would be a 3 (Effective), Administrative Controls would be a 2 (Less Effective), and PPE would be a 1 (Least Effective) PPE would be a 1 (Least Effective)), formatted exactly as defined)
    - 'Top Three Actions' (It chooses the top three actions that have the best Root Cause Identification (Utilize Causal analysis to find underlying causes, not just symptoms), Break the Causal Chain (Implement controls at multiple points; use redundancy), are Systemic and Sustainable (Focus on process, policy, and resource improvements for long-term impact), are Specific and Measurable (Define clear actions, responsibilities, and metrics for success), and reflect viable Verification and Improvement (Monitor, audit, and refine actions for ongoing effectiveness). From there, the column is populated with its ranking number and a very detailed, long, and specific (very specific to the action, including many details from the action) justification, and anything under 3 is left blank. The best corrective actions are those that address the underlying root causes of an incident through a thorough causal analysis, rather than just treating immediate symptoms. They break the chain of events at multiple points by implementing layered controls, engineering, administrative, and behavioral, to ensure redundancy and resilience. Effective corrective actions are systemic and sustainable, focusing on long-term improvements to processes, policies, and resource allocation. They are also specific, clearly defining responsibilities, timelines, and measurable outcomes, and are verified for effectiveness through ongoing monitoring and continuous improvement. This comprehensive approach ensures that corrective actions not only resolve the current issue but also prevent recurrence and strengthen organizational safety and reliability)
    - 'Assessment of Effectiveness' (For each action, what would be a good criterion to determine whether its implementation was effective?. This focuses on how effective the action is in achieving the desired end. You must explore that if the same cause reoccurred, how would the presence of the action alter the timeline of events, and what criterion would be used to evaluate its effectiveness. Clearly define the specific criteria and measurable outcomes that will be used to determine if the action is effective. Also describe the method of verification (e.g., physical testing, audit, scenario review, monitoring of performance indicators). Explain how ongoing monitoring and continuous improvement will be ensured (e.g., scheduled reviews, integration into lessons learned, feedback loops). And finally, explore how the action would disrupt the cause and effect chain in future cases, and what would be used to evaluate this disruption)
Now that you are aware of how the table should be structured, begin by creating the list of causes and actions, with their respective 'Quality of Action' rankings, and then create the table. All table rows must be sorted strictly by the 'Quality of Action' column in descending order. All actions with '5 - Most Effective' must be at the very top, followed by '4 - Highly Effective', '3 - Effective', 2 - 'Less Effective', and 1 - 'Least Effective'. You must not mix, group, or list actions by any other order, and you must not list any lower-ranked action above a higher one. Do not preserve the original order from the report; only sort by 'Quality of Action'. After presenting the tables, you must write: "All tables have been sorted by 'Quality of Action' in strict descending order." Before submitting your response, check every table and confirm that the sorting is correct. If any table is not sorted properly, you must fix it before submitting.

8. After presenting the two tables above, you must write: "All tables have been sorted by 'Quality of Action' in strict descending order." Before submitting your response, check every table and confirm that the sorting is correct. If any table is not sorted properly, you must fix it before submitting.

9. Please do not add anything extra I did not ask you to add.

10. Whenever I say to include 'all of' something. Include every single instance of what is being asked to be provided.

11. Be very detailed.

Your task:
- Analyze the content of the incident/document.
- Propose specific CORRECTIVE ACTIONS that address root causes and key risks.
- Each action must be:
  - Concrete and implementable.
  - Clearly linked to a risk or failure in the incident.
  - Framed in a professional corporate tone.

Output MUST be valid JSON with this structure:
{
  "corrective_actions": [
    {
      "id": "{{filename}}",
      "title": "...",
      "description": "...",
      "risk_addressed": "...",
      "root_cause_addressed": "...",
      "owner_suggestion": "...",
      "timeframe": "Short-term|Medium-term|Long-term",
      "impact": "High|Medium|Low",
      "confidence": 0.0
    }
  ]
}
Do not include any text outside the JSON.
"""


EVALUATION_INSTRUCTIONS = """
You are an expert reviewer of EHS corrective actions.

Your job:
- Evaluate the provided corrective actions for their potential to significantly reduce risk and improve safety.
- Consider:
  - Breadth of risk reduction.
  - Depth (severity) of issues addressed.
  - Feasibility and clarity.
  - Alignment with the customer's corrective-action guidelines.

Given a corrective-action document, extract each action and score it using this rubric:
1) RiskReduction (0-5), 
2) DowntimeAvoided (0-5), 
3) CostEffectiveness (0-5), 
4) TimeToImplement (0-5, invert score so faster=5), 
5) Repeatability (0-5).

Compute ImpactScore = 0.35*RiskReduction + 0.25*DowntimeAvoided + 0.20*CostEffectiveness + 0.10*TimeToImplement + 0.10*Repeatability.

You MUST:
- Organize in descending order the corrective actions with the most significant impact.
- Provide a short evaluation summary for each selected action.

Output MUST be valid JSON:
{
  "top_corrective_actions": [
    {
      "id": "...",
      "title": "...",
      "reason_for_selection": "...",
      "expected_impact": "High|Medium",
      "impact_score": "..."
      "comments": "...",
      "evaluation_summary": "..."
      "original_action": { ... }  // copy of original action object
    }
  ]
}
Do not include any text outside the JSON.
"""


#This is pretty standard boilerplate to wrap the agent in a ResponsesAgent() format.
class LangGraphResponsesAgent(ResponsesAgent):
    def __init__(self, agent):
        self.agent = agent

    def predict(self, request: ResponsesAgentRequest) -> ResponsesAgentResponse:
        outputs = [
            event.item
            for event in self.predict_stream(request)
            if event.type == "response.output_item.done"
        ]
        return ResponsesAgentResponse(output=outputs, custom_outputs=request.custom_inputs)

    def predict_stream(
        self,
        request: ResponsesAgentRequest,
    ) -> Generator[ResponsesAgentStreamEvent, None, None]:
        cc_msgs = to_chat_completions_input([i.model_dump() for i in request.input])

        for event in self.agent.stream({"messages": cc_msgs}, stream_mode=["updates", "messages"]):
            if event[0] == "updates":
                for node_data in event[1].values():
                    if len(node_data.get("messages", [])) > 0:
                        yield from output_to_responses_items_stream(node_data["messages"])
            # filter the streamed messages to just the generated text messages
            elif event[0] == "messages":
                try:
                    chunk = event[1][0]
                    if isinstance(chunk, AIMessageChunk) and (content := chunk.content):
                        yield ResponsesAgentStreamEvent(
                            **self.create_text_delta(delta=content, item_id=chunk.id),
                        )
                except Exception as e:
                    print(e)



#Define our tools (powerpoint and word reader)
tools = [parse_pptx_tool, parse_docx_tool]


mlflow.langchain.autolog()
agent = create_tool_calling_agent(llm, tools, EVALUATION_INSTRUCTIONS)
AGENT = LangGraphResponsesAgent(agent)
mlflow.models.set_model(AGENT)
